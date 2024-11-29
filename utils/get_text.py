"Retrievers to get LLM ready text from different sources"
import os
import logging
import re
import json
from urllib.parse import urlparse, parse_qs, urljoin
from datetime import datetime
import textwrap
import csv
import requests
from bs4 import BeautifulSoup
from pypdf import PdfReader
import docx
from atlassian import Confluence
from pptx import Presentation
import pandas as pd
from requests.auth import HTTPBasicAuth
from utils.config_utils import ConfigStore


class TxtGetterHelpers:
    """ helpers """

    @staticmethod
    def split_string(input_string):
        """ Split a string on white space or commas"""
        parts = re.split(r'[,\s]+', input_string.strip())
        parts = [item for item in parts if item != ""]
        return parts

    @staticmethod
    def get_nested_value(obj, path, default="N/A"):
        """ get a value from a nested dict using dot notation """
        keys = path.split('.')
        for key in keys:
            if obj is None:
                return default
            if isinstance(obj, dict):
                obj = obj.get(key)
            else:
                return default
        return obj if obj is not None else default

    @staticmethod
    def get_extractor_map():
        """ Get the mapping of mime types to extractor methods """

        # Map mime type to method
        extractor_map = {
            "application/pdf":
                TxtGetter.from_pdf,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                TxtGetter.from_docx,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                TxtGetter.from_pptx,
            "text/plain":
                TxtGetter.from_txt,
            "application/vnd.ms-excel":
                TxtGetter.from_xls,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                TxtGetter.from_xls,
            "text/csv":
                TxtGetter.from_csv
        }

        # Done
        return extractor_map


class TxtGetter:
    """ Class to expose methods to extract and format text from sources in an LLM ready way """

    @staticmethod
    def from_multiline_text(text):
        """ from text - nop """
        return text

    @staticmethod
    def from_pdf(file):
        """ from an adobe pdf file """
        pdf_reader = PdfReader(file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text

    @staticmethod
    def from_docx(file):
        """ from a woord doc file """
        doc = docx.Document(file)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    @staticmethod
    def from_pptx(file):
        """ from a power point file """
        prs = Presentation(file)
        text = ""


        def shape_has_table(shape):
            ''' Return true if the shape has a table '''
            try:
                # No attribute no table
                if not hasattr(shape, 'table'):
                    return False
                # Try and enumerate - will throw exception if no table
                for _i, _row in enumerate(shape.table.rows, 1):
                    pass
            except Exception as _e:
                return False

        def extract_text_from_shape(shape, indent_level=0):
            nonlocal text
            indent = "  " * indent_level

            if hasattr(shape, 'shapes'):  # Check if shape is a container
                if shape.name:
                    text += f"{indent}[Group: {shape.name}]\n"
                for sub_shape in shape.shapes:
                    extract_text_from_shape(sub_shape, indent_level + 1)
            elif shape_has_table(shape):  # Check if shape is a table
                text += f"{indent}[Table]\n"
                for i, row in enumerate(shape.table.rows, 1):
                    text += f"{indent}  Row {i}:\n"
                    for j, cell in enumerate(row.cells, 1):
                        cell_text = cell.text.strip()
                        if cell_text:
                            text += f"{indent}    Column {j}: {cell_text}\n"
            elif hasattr(shape, 'text_frame'):
                if shape.name:
                    text += f"{indent}[Shape: {shape.name}]\n"
                for i, paragraph in enumerate(shape.text_frame.paragraphs, 1):
                    if paragraph.text.strip():
                        text += f"{indent}  Paragraph {i}:\n"
                        for j, run in enumerate(paragraph.runs, 1):
                            if run.text.strip():
                                text += f"{indent}    Run {j}: {run.text.strip()}\n"
            elif hasattr(shape, 'text'):
                if shape.text.strip():
                    if shape.name:
                        text += f"{indent}[Shape: {shape.name}]: {shape.text.strip()}\n"
                    else:
                        text += f"{indent}{shape.text.strip()}\n"

        for i, slide in enumerate(prs.slides, 1):
            text += f"[Slide {i}]\n"
            for shape in slide.shapes:
                extract_text_from_shape(shape)
            text += "\n"

        return text.strip()

    @staticmethod
    def from_txt(file_path):
        """ from a text file """
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    @staticmethod
    def from_xls(file_path):
        """ from a excel file """
        df = pd.read_excel(file_path)
        return df.to_string()

    @staticmethod
    def from_csv(file_path):
        """ from a comma seperated value file """
        text = ""
        with open(file_path, 'r', encoding='utf-8') as file:
            csv_reader = csv.reader(file)
            for row in csv_reader:
                text += ", ".join(row) + "\n"
        return text


    @staticmethod
    def from_uploaded_files(uploaded_files):
        """ Extract text from uploaded files """
        total_files = len(uploaded_files)
        extracted_text = f"Text extracted from {total_files} files\n\n"

        def get_metadata(file_path):
            stat = os.stat(file_path)
            return {
                "file_name": os.path.basename(file_path),
                "file_type": file_path.split('.')[-1],
                "file_size": stat.st_size,
                "creation_date": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "last_modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            }

        def format_metadata(metadata):
            formatted = "Metadata:\n"
            for key, value in metadata.items():
                formatted += f"- {key}: {value}\n"
            return formatted

        for index, uploaded_file in enumerate(uploaded_files, 1):
            file_type = uploaded_file['type']
            file_name = uploaded_file['name']
            file_path = uploaded_file['path']

            extractor_map = TxtGetterHelpers.get_extractor_map()
            extractor = extractor_map.get(file_type)
            if not extractor:
                raise ValueError(f"Unsupported file format: {file_type}")

            metadata = get_metadata(file_path)
            extracted_text += f"File {index}/{total_files}:\n"
            extracted_text += format_metadata(metadata)

            file_content = extractor(file_path)

            # Add word count to metadata for text-based files
            if file_type not in ["application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
                word_count = len(file_content.split())
                extracted_text += f"Word count: {word_count}\n\n"

            extracted_text += f"Content:\n{file_content}\n\n"
            extracted_text += "-" * 50 + "\n\n"

        # Done
        return extracted_text

    @staticmethod
    def from_url(src):
        ''' given a url get the text '''
        response = requests.get(src)
        soup = BeautifulSoup(response.text, 'html.parser')
        text = f"Text extracted from: {src}\n\n"
        return text + ' '.join([p.get_text() for p in soup.find_all('p')])

    @staticmethod
    def from_urls(urls):
        " get text for multiple urls  "

        url_list = TxtGetterHelpers.split_string(urls)
        text = ''
        for url in url_list:
            text = text + TxtGetter.from_url(url)
            text = text + "\n\n"

        return text

    @staticmethod
    def from_jira_issue(issue_key):
        # Trim space
        issue_key = issue_key.strip()

        #  Get secrets
        jira_api_endpoint = ConfigStore.nested_get('atlassian.jira_api_endpoint')
        api_token = ConfigStore.nested_get('atlassian.api_token')
        email = ConfigStore.nested_get('atlassian.email')
        jira_url = ConfigStore.nested_get('atlassian.jira_url')

        def handle_error(response, issue_key):
            ''' Raise exception on HTTP not ok '''
            if not response.ok:
                error_msg = response.reason
                error_text = f"Could not get data for '{issue_key}' '{error_msg}'"
                raise ValueError(error_text)

        def get_issue_data(issue_key):
            url = f"{jira_url}{jira_api_endpoint}/issue/{issue_key}"
            auth = HTTPBasicAuth(email, api_token)
            headers = {"Accept": "application/json"}

            response = requests.get(url, headers=headers, auth=auth)
            handle_error(response, issue_key)
            return response.json()

        def get_issue_comments(issue_key):
            url = f"{jira_url}{jira_api_endpoint}/issue/{issue_key}/comment"
            auth = HTTPBasicAuth(email, api_token)
            headers = {"Accept": "application/json"}

            response = requests.get(url, headers=headers, auth=auth)
            handle_error(response, issue_key)
            return response.json()

        def format_description(description):
            formatted_text = ""

            def process_content(content):
                nonlocal formatted_text
                for item in content:
                    item_type = item.get('type')
                    if item_type == 'paragraph':
                        process_content(item['content'])
                    elif item_type == 'text':
                        formatted_text += item['text']
                    elif item_type == 'hardBreak':
                        formatted_text += "\n"
                    else:
                        log_msg = f"unknown item type '{item_type}' in issue '{issue_key}'"
                        logging.warning(log_msg)

            if description and 'content' in description:
                process_content(description['content'])

            return formatted_text.strip()

        def format_issue_data(issue_data, comments_data):
            get_issue = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(issue_data, path, default)
            get_comments = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(comments_data, path, default)

            formatted_output = textwrap.dedent(f"""\
                Issue Key: {get_issue('key')}
                Summary: {get_issue('fields.summary')}
                Status: {get_issue('fields.status.name')}
                Priority: {get_issue('fields.priority.name')}
                Created: {get_issue('fields.created')}
                Updated: {get_issue('fields.updated')}

                Reporter: {get_issue('fields.reporter.displayName')}

                Assignee: {get_issue('fields.assignee.displayName')}

                Description:
                {format_description(get_issue('fields.description', {}))}

                Comments:
                """)

            for comment in get_comments('comments', []):
                get_comment = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(comment, path, default)
                formatted_output += textwrap.dedent(f"""\
                    Author: {get_comment('author.displayName')}
                    Created: {get_comment('created')}
                    {format_description(get_comment('body', {}))}
                    """)

            formatted_output += "Linked Issues:\n"
            issuelinks = get_issue('fields.issuelinks', [])
            if issuelinks:
                for link in issuelinks:
                    get_link = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(link, path, default)
                    if 'outwardIssue' in link:
                        linked_issue = link['outwardIssue']
                        get_linked = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(linked_issue, path, default)
                        formatted_output += f"- {get_link('type.outward', 'Linked to')} {get_linked('key')}: {get_linked('fields.summary')}\n"
                    elif 'inwardIssue' in link:
                        linked_issue = link['inwardIssue']
                        get_linked = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(linked_issue, path, default)
                        formatted_output += f"- {get_link('type.inward', 'Linked from')} {get_linked('key')}: {get_linked('fields.summary')}\n"
            else:
                formatted_output += "No linked issues found.\n"

            return formatted_output.strip()

        # Main execution
        issue_key = issue_key.strip()
        issue_data = get_issue_data(issue_key)
        comments_data = get_issue_comments(issue_key)
        formatted_issue = format_issue_data(issue_data, comments_data)
        return formatted_issue

    @staticmethod
    def from_jira_issues(issue_keys):
        " get text for multiple jira tickets "

        issues_keys_list = TxtGetterHelpers.split_string(issue_keys)
        text = ''
        for issue_key in issues_keys_list:
            text = text + TxtGetter.from_jira_issue(issue_key)
            text = text + "\n\n"

        return text

    @staticmethod
    def from_jql_query(jql_query, page_size=50, max_results=100):
        """ Get text from the results of a JQL query """

        def handle_error(response):
            if not response.ok:
                error_msg = response.reason
                error_text = f"Error executing JQL query: {error_msg}"
                raise ValueError(error_text)

        def get_issues_data(jql, start_at=0, max_results=50):

            # Get config
            jira_url = ConfigStore.nested_get('atlassian.jira_url')
            jira_api_endpoint = ConfigStore.nested_get('atlassian.jira_api_endpoint')
            api_token = ConfigStore.nested_get('atlassian.api_token')
            email = ConfigStore.nested_get('atlassian.email')

            # Format
            url = f"{jira_url}{jira_api_endpoint}/search"

            auth = HTTPBasicAuth(email, api_token)
            headers = {
                "Accept": "application/json",
                "Content-Type": "application/json"
            }

            payload = json.dumps({
                "jql": jql,
                "startAt": start_at,
                "maxResults": max_results,
                "fields": [
                    "summary", "status", "priority", "created", "updated",
                    "reporter", "assignee", "description", "comment", "issuelinks"
                ]
            })

            response = requests.post(url, data=payload, headers=headers, auth=auth)
            handle_error(response)
            return response.json()

        def format_description(description):
            formatted_text = ""

            def process_content(content):
                nonlocal formatted_text
                for item in content:
                    item_type = item.get('type')
                    if item_type == 'paragraph':
                        process_content(item['content'])
                    elif item_type == 'text':
                        formatted_text += item['text']
                    elif item_type == 'hardBreak':
                        formatted_text += "\n"
                    else:
                        log_msg = f"unknown item type '{item_type}' in JQL result '{jql_query}'"
                        logging.warning(log_msg)

            if description and 'content' in description:
                process_content(description['content'])

            return formatted_text.strip()

        def format_issues_data(issues_data):
            formatted_output = ""

            for issue in issues_data['issues']:
                get_issue = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(issue, path, default)

                issue_text = textwrap.dedent(f"""\
                    Issue Key: {get_issue('key')}
                    Summary: {get_issue('fields.summary')}
                    Status: {get_issue('fields.status.name')}
                    Priority: {get_issue('fields.priority.name')}
                    Created: {get_issue('fields.created')}
                    Updated: {get_issue('fields.updated')}
                    Reporter: {get_issue('fields.reporter.displayName')}
                    Assignee: {get_issue('fields.assignee.displayName')}

                    Description:
                    {format_description(get_issue('fields.description', {}))}

                    Comments:
                    """)
                formatted_output += issue_text

                for comment in get_issue('fields.comment.comments', []):
                    get_comment = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(comment, path, default)
                    comment_text = textwrap.dedent(f"""\
                        Author: {get_comment('author.displayName')}
                        Created: {get_comment('created')}
                        {format_description(get_comment('body', {}))}
                        """)
                    formatted_output += comment_text

                formatted_output += "Linked Issues:\n"
                for link in get_issue('fields.issuelinks', []):
                    get_link = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(link, path, default)
                    if 'outwardIssue' in link:
                        linked_issue = link['outwardIssue']
                        get_linked = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(linked_issue, path, default)
                        formatted_output += f"- {get_link('type.outward', 'Linked to')} {get_linked('key')}: {get_linked('fields.summary')}\n"
                    elif 'inwardIssue' in link:
                        linked_issue = link['inwardIssue']
                        get_linked = lambda path, default="N/A": TxtGetterHelpers.get_nested_value(linked_issue, path, default)
                        formatted_output += f"- {get_link('type.inward', 'Linked from')} {get_linked('key')}: {get_linked('fields.summary')}\n"

                formatted_output += "\n---\n"  # Separator between issues

            return formatted_output.strip()

        # Main execution with pagination
        all_issues = []
        start_at = 0

        while len(all_issues) < max_results:
            issues_data = get_issues_data(jql_query, start_at, page_size)
            all_issues.extend(issues_data['issues'])

            if len(issues_data['issues']) < page_size or len(all_issues) >= issues_data['total']:
                break

            start_at += page_size

        # Trim the results to max_results if necessary
        all_issues = all_issues[:max_results]

        # Create a new dictionary with the structure expected by format_issues_data
        formatted_data = {'issues': all_issues}

        formatted_issues = format_issues_data(formatted_data)
        return formatted_issues

    @staticmethod
    def from_confluence_page(page_url_or_id):
        """Extract text and metadata from a Confluence page, including tables and embedded components."""

        class ConfluencePageExtractor:
            def __init__(self, url, username, api_token):
                self.confluence = Confluence(
                    url=url,
                    username=username,
                    password=api_token
                )

            def extract_page_id_from_url(self, url):
                parsed_url = urlparse(url)
                path = parsed_url.path

                if '/pages/' in path:
                    page_id = path.split('/pages/')[1].split('/')[0]
                    return page_id

                query_params = parse_qs(parsed_url.query)
                if 'pageId' in query_params:
                    return query_params['pageId'][0]

                raise ValueError(f"Unable to extract page ID from '{url}'")

            def from_confluence_page(self, page):
                page_id = self.extract_page_id_from_url(page) if page.startswith('http') else page

                # Fetch the page content with the 'view' representation
                page_content = self.confluence.get_page_by_id(page_id, expand='body.view,version,metadata.labels')

                html_content = page_content['body']['view']['value']
                soup = BeautifulSoup(html_content, 'html.parser')

                # Extract metadata
                title = page_content['title']
                author = page_content['version']['by']['displayName']
                last_updated = datetime.fromisoformat(page_content['version']['when'].rstrip('Z')).strftime('%Y-%m-%d %H:%M:%S')
                labels = [label['name'] for label in page_content['metadata']['labels']['results']]

                # Extract links
                links = [{'text': a.text, 'href': urljoin(page, a['href'])} for a in soup.find_all('a', href=True)]

                # Extract text, preserving structure and including rendered components
                paragraphs = []
                for element in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'li', 'table', 'div']):
                    if element.name == 'table':
                        paragraphs.append("\n\nTABLE:\n" + element.get_text(separator=' | ', strip=True) + "\n")
                    elif element.name == 'div' and 'data-macro-name' in element.attrs:
                        macro_name = element['data-macro-name']
                        paragraphs.append(f"\n[{macro_name.upper()}]\n{element.get_text(separator=' ', strip=True)}\n")
                    else:
                        text = element.get_text(strip=True)
                        if text:
                            if element.name.startswith('h'):
                                paragraphs.append(f"\n\n{element.name.upper()}: {text}\n")
                            elif element.name == 'li':
                                paragraphs.append(f"- {text}")
                            else:
                                paragraphs.append(text)

                text = " ".join(paragraphs)
                text = re.sub(r'\s+', ' ', text)
                text = re.sub(r'\n{3,}', '\n\n', text)

                # Compile the final output using textwrap.dedent
                output = textwrap.dedent(f"""\
                    Title: {title}
                    Author: {author}
                    Last Updated: {last_updated}
                    Labels: {', '.join(labels)}
                    URL: {page}

                    Content:
                    {text}

                    Links:
                    {json.dumps(links, indent=2)}
                """).strip()

                return output

        # Create extractor
        extractor = ConfluencePageExtractor(
            url=ConfigStore.nested_get('atlassian.jira_url'),
            username=ConfigStore.nested_get('atlassian.email'),
            api_token=ConfigStore.nested_get('atlassian.api_token')
        )
        content = extractor.from_confluence_page(page_url_or_id)
        return content

    @staticmethod
    def from_confluence_pages(page_urls_or_ids):
        """ Get text from a list of confluence pages """

        page_list = TxtGetterHelpers.split_string(page_urls_or_ids)
        text = ''
        for page in page_list:
            text = text + TxtGetter.from_confluence_page(page)
            text = text + "\n\n"

        return text
